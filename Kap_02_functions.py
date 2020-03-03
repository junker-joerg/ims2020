# Tabelle anlegen
# ! Achtung: alles in pandas machen
# ToDo  zur Installation: pip install python-pptx
# ! kommt vom KPMG Rechner ... nochmal

"""
slide_layouts[x]
0 Title (presentation title slide)
1 Title and Content
2 Section Header (sometimes called Segue)
3 Two Content (side by side bullet textboxes)
4 Comparison (same but additional title for each side by side content box)
5 Title Only
6 Blank
7 Content with Caption
8 Picture with Caption
 
 shapes - 6 verschiedene "Hauptarten"
 1 autoshape: 180 verschiedene
 2 picturepptx
 3 grahic frame - kann Tabelle, Diagramm
 4 group shape - eine Menge von Shapes
 5 line / contector
 6 content part

Gebräuchliste shapes:
shape shapes – auto shapes with fill and an outline
text boxes – auto shapes with no fill and no outline
placeholders – auto shapes that can appear on a slide layout or master and be inherited on slides that use that layout, allowing content to be added that takes on the formatting of the placeholder
line/connector – as described above
picture – as described above
table – that row and column thing
chart – pie chart, line chart, etc.
smart art – not supported yet, although preserved if present
media clip – video or audio

Alles Shapes werden über einen Shape-Tree verwaltet - Zugriff: 

    shapes = slide.shape


 """

from pptx import Presentation
from pptx.util import Inches
import os

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

shapes.title.text = 'Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Martin'
table.cell(0, 1).text = 'Köhler'

# write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'

prs.save('test.pptx')
os.startfile('test.pptx') # https://pythonprogramming.altervista.org/inserting-an-image-in-powerpoint-with-python/
